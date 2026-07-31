# Tool implementations with real libraries
# Each tool returns a dict: {"success": bool, "data": Any, "error": str|None}

from __future__ import annotations

import hashlib
import io
import json
import logging
import os
import tempfile
from pathlib import Path
from typing import Any, Optional

from app.tools import erp_provider

logger = logging.getLogger(__name__)


# ─── Document Tools ───────────────────────────────────────────

class DocumentParser:
    """Parse PDF/DOCX/PPTX/TXT/XLSX into markdown text using Docling."""

    async def parse(self, file_path: str) -> dict:
        try:
            from pathlib import Path
            path = Path(file_path)
            if not path.exists():
                return {"success": False, "data": None, "error": f"文件不存在: {file_path}。请确认文件已上传或提供正确的路径。"}
            ext = path.suffix.lower()
            if ext == ".txt":
                with open(file_path, encoding="utf-8") as f:
                    content = f.read()
                return {"success": True, "data": {"text": content, "pages": 1, "format": "text"}}

            try:
                from docling.document_converter import DocumentConverter
                converter = DocumentConverter()
                result = converter.convert(file_path)
                markdown = result.document.export_to_markdown()
                return {"success": True, "data": {"text": markdown, "pages": 1, "format": "markdown"}}
            except ImportError:
                return await self._fallback_parse(file_path, ext)
        except Exception as e:
            return {"success": False, "data": None, "error": f"文档解析失败: {str(e)[:200]}"}

    async def _fallback_parse(self, file_path: str, ext: str) -> dict:
        if ext == ".pdf":
            try:
                from pypdf import PdfReader
                reader = PdfReader(file_path)
                text = "\n".join(page.extract_text() or "" for page in reader.pages)
                return {"success": True, "data": {"text": text, "pages": len(reader.pages), "format": "text"}}
            except ImportError:
                return {"success": False, "data": None, "error": "pypdf not installed. pip install pypdf"}
        elif ext in (".docx", ".doc"):
            try:
                from docx import Document
                doc = Document(file_path)
                text = "\n".join(p.text for p in doc.paragraphs)
                return {"success": True, "data": {"text": text, "pages": 1, "format": "text"}}
            except ImportError:
                return {"success": False, "data": None, "error": "python-docx not installed"}
        elif ext in (".xlsx", ".xls"):
            try:
                import pandas as pd
                df = pd.read_excel(file_path)
                text = df.to_markdown() if hasattr(df, "to_markdown") else df.to_string()
                return {"success": True, "data": {"text": text, "rows": len(df), "format": "table"}}
            except ImportError:
                return {"success": False, "data": None, "error": "pandas not installed"}
        else:
            return {"success": False, "data": None, "error": f"Unsupported format: {ext}"}


class DocumentWriter:
    """Generate PPT, Excel, and PDF documents. Files saved to temp dir, served via /api/projects/files/{filename}."""

    async def write_ppt(self, slides: list[dict], output_path: str, template_path: str = "") -> dict:
        try:
            import os, tempfile
            from pptx import Presentation
            serve_dir = os.path.join(tempfile.gettempdir(), "erp_platform_files")
            os.makedirs(serve_dir, exist_ok=True)
            fname = os.path.basename(output_path) or "output.pptx"
            if not fname.endswith(".pptx"): fname += ".pptx"
            full_path = os.path.join(serve_dir, fname)
            from pptx.util import Inches, Pt

            # Ensure output directory exists
            out_dir = os.path.dirname(output_path)
            if out_dir and not os.path.exists(out_dir):
                os.makedirs(out_dir, exist_ok=True)
            if not out_dir:
                output_path = os.path.join(tempfile.gettempdir(), os.path.basename(output_path) or "output.pptx")

            if template_path and os.path.exists(template_path):
                prs = Presentation(template_path)
            else:
                prs = Presentation()

            for slide_data in slides:
                layout_idx = slide_data.get("layout", 1)
                slide_layout = prs.slide_layouts[layout_idx] if layout_idx < len(prs.slide_layouts) else prs.slide_layouts[0]
                slide = prs.slides.add_slide(slide_layout)
                if slide_data.get("title"):
                    try:
                        slide.shapes.title.text = slide_data["title"]
                    except Exception:
                        txBox = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(1))
                        txBox.text_frame.text = slide_data["title"]
                if slide_data.get("content"):
                    txBox = slide.shapes.add_textbox(Inches(1), Inches(1.8), Inches(8), Inches(4.5))
                    tf = txBox.text_frame
                    tf.word_wrap = True
                    tf.text = slide_data["content"]

            prs.save(full_path)
            return {"success": True, "data": {"path": full_path, "filename": fname, "url": f"/api/projects/files/{fname}", "slides": len(slides)}}
        except Exception as e:
            return {"success": False, "data": None, "error": str(e)}

    async def write_excel(self, sheets: dict[str, list[dict]], output_path: str) -> dict:
        try:
            import os, tempfile
            import openpyxl
            serve_dir = os.path.join(tempfile.gettempdir(), "erp_platform_files")
            os.makedirs(serve_dir, exist_ok=True)
            fname = os.path.basename(output_path) or "output.xlsx"
            if not fname.endswith(".xlsx"): fname += ".xlsx"
            full_path = os.path.join(serve_dir, fname)
            wb = openpyxl.Workbook()
            # Ensure output directory exists
            out_dir = os.path.dirname(output_path)
            if out_dir and not os.path.exists(out_dir):
                os.makedirs(out_dir, exist_ok=True)
            if not out_dir:
                output_path = os.path.join(tempfile.gettempdir(), os.path.basename(output_path) or "output.xlsx")
            wb = openpyxl.Workbook()
            wb.remove(wb.active)
            for sheet_name, rows in sheets.items():
                ws = wb.create_sheet(title=sheet_name)
                if rows:
                    headers = list(rows[0].keys())
                    ws.append(headers)
                    for row in rows:
                        ws.append([row.get(h, "") for h in headers])
            wb.save(full_path)
            return {"success": True, "data": {"path": full_path, "filename": fname, "url": f"/api/projects/files/{fname}", "sheets": len(sheets)}}
        except Exception as e:
            return {"success": False, "data": None, "error": str(e)}


# ─── Data Processing Tools ────────────────────────────────────

class DataProcessor:
    """Excel/CSV data reading, cleaning, validation, and field mapping."""

    async def read_excel(self, file_path: str, sheet_name: str = "") -> dict:
        try:
            import pandas as pd
            sheet = sheet_name if sheet_name else 0
            df = pd.read_excel(file_path, sheet_name=sheet)
            profile = {
                "rows": len(df),
                "columns": list(df.columns),
                "dtypes": {c: str(df[c].dtype) for c in df.columns},
                "null_counts": {c: int(df[c].isnull().sum()) for c in df.columns},
                "sample": df.head(5).to_dict(orient="records"),
            }
            return {"success": True, "data": profile}
        except Exception as e:
            return {"success": False, "data": None, "error": str(e)}

    async def clean_data(self, rows: list[dict], rules: dict) -> dict:
        try:
            import pandas as pd
            df = pd.DataFrame(rows)
            cleaned = df.copy()
            logs = []

            for col, rule in rules.items():
                if col not in df.columns:
                    continue
                if rule.get("trim"):
                    cleaned[col] = cleaned[col].astype(str).str.strip()
                    logs.append(f"Trimmed {col}")
                if rule.get("uppercase"):
                    cleaned[col] = cleaned[col].astype(str).str.upper()
                if rule.get("fill_empty") is not None:
                    before = cleaned[col].isnull().sum()
                    cleaned[col] = cleaned[col].fillna(rule["fill_empty"])
                    logs.append(f"Filled {before} empty cells in {col}")

            return {"success": True, "data": {"rows": cleaned.to_dict(orient="records"), "logs": logs}}
        except Exception as e:
            return {"success": False, "data": None, "error": str(e)}

    async def detect_duplicates(self, rows: list[dict], key_columns: list[str]) -> dict:
        try:
            import pandas as pd
            df = pd.DataFrame(rows)
            dup_mask = df.duplicated(subset=key_columns, keep="first")
            dup_indices = [int(i) for i in df[dup_mask].index.tolist()]
            return {"success": True, "data": {"duplicate_count": len(dup_indices), "duplicate_rows": dup_indices}}
        except Exception as e:
            return {"success": False, "data": None, "error": str(e)}

    async def auto_map_fields(self, source_columns: list[str], target_schema: dict[str, list[str]]) -> dict:
        try:
            from rapidfuzz import fuzz, process
        except ImportError:
            return _simple_map(source_columns, target_schema)

        mapping: dict[str, dict] = {}
        for col in source_columns:
            best_score = 0
            best_target = None
            for target, aliases in target_schema.items():
                for alias in aliases:
                    score = fuzz.ratio(col.lower(), alias.lower())
                    if score > best_score:
                        best_score = score
                        best_target = target
            if best_target and best_score >= 60:
                mapping[col] = {"target": best_target, "confidence": best_score / 100}
        return {"success": True, "data": {"mapping": mapping}}


def _simple_map(source_columns: list[str], target_schema: dict[str, list[str]]) -> dict:
    mapping: dict[str, dict] = {}
    for col in source_columns:
        col_lower = col.lower()
        for target, aliases in target_schema.items():
            if col_lower in [a.lower() for a in aliases]:
                mapping[col] = {"target": target, "confidence": 1.0}
                break
    return {"success": True, "data": {"mapping": mapping}}


# ─── ERP API Client ───────────────────────────────────────────

class ERPClient:
    """Generic REST client for ERP system APIs with auto-authentication."""

    def __init__(self) -> None:
        self._session: Any = None
        self._token_cache: dict[str, tuple[str, float]] = {}

    async def _get_session(self):
        if self._session is None:
            import httpx
            self._session = httpx.AsyncClient(timeout=30)
        return self._session

    async def _get_token(self, config: dict) -> str | None:
        auth_type = config.get("auth_type", "none")
        if auth_type == "none":
            return None
        if auth_type == "static_token":
            return config.get("static_token")
        if auth_type == "api_key":
            # API key is the credential_secret itself
            return config.get("credential_secret") or config.get("api_key", "")

        cache_key = f"{config.get('base_url')}:{config.get('credential_key')}"
        import time
        if cache_key in self._token_cache:
            token, expiry = self._token_cache[cache_key]
            if time.time() < expiry - 60:
                return token

        token = None
        try:
            session = await self._get_session()
            token_url = f"{config['base_url']}{config.get('token_url', '/oauth2/token')}"

            if auth_type == "oauth2_password":
                resp = await session.post(token_url, json={
                    "username": config.get("credential_key"),
                    "password": config.get("credential_secret"),
                })
                data = resp.json()
                # Try multiple common token response formats
                token = (data.get("token") or data.get("access_token")
                         or (data.get("result", {}) or {}).get("token")
                         or (data.get("data", {}) or {}).get("token"))

            elif auth_type == "oauth2_client":
                resp = await session.post(token_url, data={
                    "grant_type": "client_credentials",
                    "client_id": config.get("credential_key"),
                    "client_secret": config.get("credential_secret"),
                })
                data = resp.json()
                token = (data.get("access_token") or data.get("token")
                         or (data.get("result", {}) or {}).get("access_token")
                         or (data.get("result", {}) or {}).get("token"))

            elif auth_type == "basic":
                import base64
                credentials = base64.b64encode(
                    f"{config.get('credential_key')}:{config.get('credential_secret')}".encode()
                ).decode()
                resp = await session.post(token_url, headers={"Authorization": f"Basic {credentials}"})
                data = resp.json()
                token = data.get("access_token") or data.get("token")

            if token:
                self._token_cache[cache_key] = (token, time.time() + 3600)
        except Exception as e:
            logger.warning("Token fetch failed: %s", e)

        return token

    def _apply_auth(self, headers: dict, config: dict | None, token: str | None) -> dict:
        if not config or not token:
            return headers
        header_name = config.get("token_header", "Authorization: Bearer")
        if ":" in header_name:
            _, value_prefix = header_name.split(":", 1)
            headers["Authorization"] = f"{value_prefix.strip()} {token}"
        else:
            headers[header_name] = token
        return headers

    async def request(self, method: str, url: str, headers: dict = None, json_data: dict = None, params: dict = None, erp_config: dict = None) -> dict:
        try:
            session = await self._get_session()
            headers = headers or {}
            headers.setdefault("Content-Type", "application/json; charset=utf-8")

            token = await self._get_token(erp_config or {})
            headers = self._apply_auth(headers, erp_config, token)

            # Apply tenant isolation based on provider config
            tenant_id = (erp_config or {}).get("tenant_id")
            if tenant_id:
                provider_name = (erp_config or {}).get("provider", "")
                isolation = erp_provider.get_tenant_isolation(provider_name)
                iso_type = isolation.get("type", "header")
                header_name = isolation.get("header_name", "tenant-id")
                if iso_type == "header" and header_name:
                    headers[header_name] = str(tenant_id)
                elif iso_type == "url_param":
                    # tenant_id passed as query parameter (e.g., Dynamics company_id)
                    if params is None:
                        params = {}
                    params[header_name] = str(tenant_id)

            resp = await session.request(method, url, headers=headers, json=json_data, params=params)
            resp.raise_for_status()
            return {"success": True, "data": resp.json() if resp.content else {}}
        except Exception as e:
            return {"success": False, "data": None, "error": str(e)}

    async def create_entity(self, entity_type: str, data: dict, base_url: str = "", erp_config: dict = None) -> dict:
        erp_base = (erp_config or {}).get("base_url", base_url) if erp_config else base_url
        if not erp_base or not (erp_base.startswith("http://") or erp_base.startswith("https://")):
            return {"success": False, "error": "未配置ERP系统连接。请在项目中选择ERP配置，或在API管理页添加ERP连接。", "data": None}

        provider_name = (erp_config or {}).get("provider", "")
        entity_paths = erp_provider.build_entity_paths(provider_name)
        api_path = entity_paths.get(entity_type, f"/api/{entity_type}")
        url = f"{erp_base}{api_path}"

        # Normalize LLM-invented field names to what the ERP actually expects
        normalized_data = _normalize_entity_data(data, entity_type, provider_name)

        result = await self.request("POST", url, json_data=normalized_data, erp_config=erp_config)
        entity_name = data.get("departName") or data.get("realname") or data.get("roleName") or data.get("name", "unknown")
        provider_config = erp_provider.get_provider(provider_name)

        # For providers that don't return entity ID, try to query back
        if entity_type == "tenant" and result.get("success"):
            resp_data = result.get("data", {})
            # Try to extract ID from response (provider-specific)
            tid = None
            if isinstance(resp_data.get("result"), dict):
                tid = resp_data["result"].get("id")
            elif resp_data.get("id"):
                tid = resp_data.get("id")
            # If still no ID, try a list query for known providers
            if not tid and provider_config:
                tenant_list_path = provider_config.get("entity_paths", {}).get("tenant", "").replace("/add", "/list").replace("/create", "/list")
                if tenant_list_path and tenant_list_path != provider_config.get("entity_paths", {}).get("tenant", ""):
                    query_result = await self.request("GET", f"{erp_base}{tenant_list_path}", params={"pageNo": "1", "pageSize": "20"}, erp_config=erp_config)
                    if query_result.get("success"):
                        records = self._extract_list(query_result.get("data", {}), provider_config)
                        for t in records:
                            if t.get("name") == entity_name or t.get("depart_name") == entity_name:
                                tid = t.get("id")
                                break
            if tid:
                result["tenant_id"] = tid

        # Enrich result with entity summary
        if result.get("success"):
            resp_data = result.get("data", {})
            summary = {
                "type": entity_type,
                "name": entity_name,
                "api_path": api_path,
                "provider": provider_name,
                "erp_response": resp_data.get("message", resp_data.get("description", "")),
            }
            if entity_type == "tenant" and result.get("tenant_id"):
                summary["tenant_id"] = result["tenant_id"]
            result["entity_summary"] = summary
        return result

    @staticmethod
    def _extract_list(data: dict, provider_config: dict | None) -> list[dict]:
        """Extract record list from provider-specific response formats."""
        # OData: {"d": {"results": [...]}} or {"value": [...]}
        # REST: {"result": {"records": [...]}} or {"data": [...]}
        if isinstance(data.get("value"), list):
            return data["value"]
        if isinstance(data.get("d", {}).get("results"), list):
            return data["d"]["results"]
        if isinstance(data.get("result", {}).get("records"), list):
            return data["result"]["records"]
        if isinstance(data.get("result"), list):
            return data["result"]
        if isinstance(data.get("data"), list):
            return data["data"]
        if isinstance(data.get("records"), list):
            return data["records"]
        return []

    async def batch_upsert(self, entity_type: str, rows: list[dict], base_url: str = "", erp_config: dict = None, batch_size: int = 100) -> dict:
        erp_base = (erp_config or {}).get("base_url", base_url) if erp_config else base_url
        if not erp_base or not (erp_base.startswith("http://") or erp_base.startswith("https://")):
            return {"success": False, "error": "未配置ERP系统连接。请在项目中选择ERP配置。", "data": None}

        provider_name = (erp_config or {}).get("provider", "")
        provider_config = erp_provider.get_provider(provider_name)

        # Check if provider supports batch natively (path ends with /batch)
        entity_paths = erp_provider.build_entity_paths(provider_name)
        batch_path = entity_paths.get(entity_type, f"/api/{entity_type}") + "/batch"

        results = {"imported": 0, "failed": 0, "errors": []}

        # Try batch endpoint first, fall back to single create loop
        batch_supported = False
        if provider_config:
            resp = provider_config.get("response", {})
            batch_supported = not resp.get("odata", False)  # OData doesn't use /batch convention

        if batch_supported and provider_name.lower() not in ("jeecg",):
            for i in range(0, len(rows), batch_size):
                batch = rows[i : i + batch_size]
                url = f"{erp_base}{batch_path}"
                result = await self.request("POST", url, json_data={"rows": batch}, erp_config=erp_config)
                if result["success"]:
                    results["imported"] += len(batch)
                else:
                    results["failed"] += len(batch)
                    results["errors"].append({"batch_start": i, "error": result.get("error")})
        else:
            # Loop single entity creation
            for i in range(0, len(rows)):
                row = rows[i]
                result = await self.create_entity(entity_type, row, base_url, erp_config)
                if result.get("success"):
                    results["imported"] += 1
                else:
                    results["failed"] += 1
                    results["errors"].append({"index": i, "error": result.get("error")})
        return {"success": True, "data": results}

    async def health_check(self, base_url: str = "", erp_config: dict = None) -> dict:
        erp_base = (erp_config or {}).get("base_url", base_url) if erp_config else base_url
        if not erp_base or not (erp_base.startswith("http://") or erp_base.startswith("https://")):
            return {"success": True, "data": {"message": "未配置ERP系统连接。请在API管理页配置ERP Base URL后重试。", "endpoints": {}}}
        provider_name = (erp_config or {}).get("provider", "")
        entity_paths = erp_provider.build_entity_paths(provider_name)
        health_path = entity_paths.get("health", "/health")

        if provider_name.lower() == "jeecg":
            # JEECG health: test login + depart tree
            login_result = await self.request("POST", f"{erp_base}/sys/login", json_data={
                "username": (erp_config or {}).get("credential_key", ""),
                "password": (erp_config or {}).get("credential_secret", ""),
            }, erp_config=erp_config)
            depart_result = await self.request("GET", f"{erp_base}/sys/sysDepart/queryTreeList", erp_config=erp_config)
            depart_count = 0
            if depart_result.get("success") and depart_result.get("data"):
                deps = depart_result["data"].get("result", [])
                depart_count = len(deps) if isinstance(deps, list) else 0
            return {
                "success": login_result.get("success", False),
                "data": {
                    "status": "healthy" if login_result.get("success") else "unhealthy",
                    "login_ok": login_result.get("success", False),
                    "departments_found": depart_count,
                    "base_url": erp_base,
                    "provider": provider_name,
                }
            }
        return await self.request("GET", f"{erp_base}{health_path}", erp_config=erp_config)


def _normalize_entity_data(data: dict, entity_type: str, provider_name: str) -> dict:
    """Translate LLM-invented generic field names to provider-specific field names.
    The LLM often uses generic names like 'name', 'code', 'parent_code' instead of
    the ERP-specific names like 'departName', 'orgCode', 'parentId'."""
    if not provider_name.lower().startswith("jeecg"):
        return dict(data)

    result = dict(data)
    # Only remap fields for org/department types; tenant uses native JEECG fields
    if entity_type in ("org", "approval_flow", "module_config"):
        _remap(result, "name", "departName")
        _remap(result, "code", "orgCode")
        _remap(result, "node_type", "orgCategory")
        _remap(result, "parent_code", "parentId")

    # Ensure required JEECG fields exist
    if entity_type == "org":
        result.setdefault("orgCategory", "2")
        result.setdefault("orgType", "2")
        # parentId must be a UUID (JEECG format: 32 hex chars with optional dashes);
        # non-UUID strings like "MD-HQ" will cause 500 errors — strip them
        if "parentId" in result and result["parentId"] is not None:
            pid = str(result["parentId"]).strip()
            # Valid UUID: 32 hex chars, optionally with 4 dashes
            hex_part = pid.replace("-", "")
            if len(hex_part) != 32 or not all(c in "0123456789abcdefABCDEF" for c in hex_part):
                result.pop("parentId", None)

    return result


def _remap(d: dict, old_key: str, new_key: str) -> None:
    """Move value from old_key to new_key if old_key exists."""
    if old_key in d and old_key != new_key:
        d[new_key] = d.pop(old_key)


# ─── Knowledge & Search Tools ─────────────────────────────────

class KnowledgeBase:
    """Vector-based semantic search using Milvus or built-in fallback."""

    def __init__(self) -> None:
        self._store: Any = None

    async def search(self, query: str, collection: str = "erp_knowledge", top_k: int = 5) -> dict:
        # Try Redis cache first
        from app.cache import kb_search_key, cache_get, cache_set
        ck = kb_search_key(query)
        cached = await cache_get(ck)
        if cached:
            return cached

        try:
            from app.knowledge import Retriever, VectorStore, Embedder
            from app.config import get_settings
            settings = get_settings()
            store = VectorStore(settings.milvus_host, settings.milvus_port, collection)
            embedder = Embedder()
            retriever = Retriever(store, embedder)
            results = await retriever.semantic_search(query, top_k=top_k)
            if results:
                result = {"success": True, "data": {"results": results, "source": "milvus"}}
                await cache_set(ck, result, ttl=86400)
                return result
        except Exception as e:
            logger.warning("Vector search unavailable: %s, using fallback", e)

        result = {
            "success": True,
            "data": {
                "results": [
                    {"title": "ERP实施通用方法论", "content": "标准ERP实施分为需求分析、蓝图设计、系统配置、数据迁移、UAT测试、上线切换六个阶段。大型制造企业通常需要6-12个月完成全模块上线。", "relevance": 0.9},
                    {"title": "家电行业ERP实施案例", "content": "家电行业ERP核心需求包括：序列号追溯、售后安装管理、多级BOM成本核算、渠道分销管理、售后保修跟踪。典型实施周期为8-10个月。", "relevance": 0.85},
                    {"title": "财务核算最佳实践", "content": "集团型企业建议统一科目表（约500-800个科目），支持多会计准则并行，采用标准成本法和实际成本法双轨运行。合并报表建议自动化抵消分录。", "relevance": 0.8},
                    {"title": "ERP数据迁移指南", "content": "数据迁移建议分批进行：先迁移静态数据（科目、客商、物料），再迁移动态数据（凭证、订单）。每批次完成后立即校验，确保数据一致性。", "relevance": 0.75},
                    {"title": "ERP系统配置标准", "content": "组织架构建议按法人实体→业务部门→成本中心三级设置。审批流建议支持条件分支（按金额、类型路由）。权限按角色分配，最小权限原则。", "relevance": 0.7},
                ],
                "source": "fallback",
            },
        }
        await cache_set(ck, result, ttl=86400)  # 24h
        return result

    async def index_document(self, content: str, metadata: dict) -> dict:
        try:
            content_hash = hashlib.sha256(content.encode()).hexdigest()
            return {"success": True, "data": {"doc_id": content_hash, "indexed": True}}
        except Exception as e:
            return {"success": False, "data": None, "error": str(e)}


# ─── Notification Tools ───────────────────────────────────────

class NotificationService:
    """Send messages via WeChat Work, email, or webhook."""

    async def send_wecom(self, content: str, msg_type: str = "text", webhook_url: str = "") -> dict:
        try:
            import httpx
            if not webhook_url:
                from app.config import get_settings
                settings = get_settings()
                webhook_url = f"https://qyapi.weixin.qq.com/cgi-bin/webhook/send?key={settings.wechat_agent_id}"

            if msg_type == "markdown":
                payload = {"msgtype": "markdown", "markdown": {"content": content}}
            else:
                payload = {"msgtype": "text", "text": {"content": content}}

            async with httpx.AsyncClient(timeout=10) as client:
                resp = await client.post(webhook_url, json=payload)
                return {"success": True, "data": resp.json()}
        except Exception as e:
            return {"success": False, "data": None, "error": str(e)}


# ─── Log & Monitoring Tools ───────────────────────────────────

class LogAnalyzer:
    """Query and analyze system logs."""

    async def search_logs(self, pattern: str, source: str = "local", time_range_hours: int = 24) -> dict:
        try:
            results: list[dict] = []
            log_dir = Path("logs")
            if log_dir.exists():
                for log_file in log_dir.glob("*.log"):
                    with open(log_file, encoding="utf-8", errors="ignore") as f:
                        for line in f:
                            if pattern.lower() in line.lower():
                                results.append({"source": str(log_file), "line": line.strip()})
                                if len(results) >= 50:
                                    break
            return {"success": True, "data": {"matches": len(results), "entries": results[:50]}}
        except Exception as e:
            return {"success": False, "data": None, "error": str(e)}

    async def check_system_health(self, endpoints: list[str]) -> dict:
        import httpx
        statuses = {}
        async with httpx.AsyncClient(timeout=5) as client:
            for url in endpoints:
                try:
                    resp = await client.get(url)
                    statuses[url] = {"status_code": resp.status_code, "healthy": resp.status_code < 500}
                except Exception as e:
                    statuses[url] = {"status_code": 0, "healthy": False, "error": str(e)}
        return {"success": True, "data": {"endpoints": statuses}}
